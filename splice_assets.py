"""splice_assets.py — replace ASSET_PLACEHOLDER shapes with real images.

Walks a .pptx, finds each shape whose name starts with `ASSET_PLACEHOLDER:`,
parses the asset_id + fit_mode, looks up the binary in the asset folder,
and inserts it at the placeholder's geometry.

Usage:
  python splice_assets.py in.pptx --assets /path/to/assets -o out.pptx

Missing assets and unsupported types (SVG/XML) leave the placeholder in place
and append to <out>.warnings.txt.
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

import yaml
from PIL import Image
from pptx import Presentation
from pptx.util import Emu

PLACEHOLDER_PREFIX = "ASSET_PLACEHOLDER:"

_SUPPORTED_RASTER = {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff", ".webp"}


def _read_asset_index(asset_dir: Path) -> dict[str, dict]:
    """Build id → sidecar dict from *.yaml files in asset_dir.

    Each entry also has `_abs_path` for the binary, resolved from the
    yaml's `file:` field (or by matching the yaml stem to a binary).
    """
    index = {}
    for y in asset_dir.glob("*.yaml"):
        if y.name in {"theme.yaml", "asset_tag_vocab.yaml"}:
            continue
        try:
            with open(y) as f:
                data = yaml.safe_load(f) or {}
        except yaml.YAMLError:
            continue
        asset_id = data.get("id") or y.stem
        file_field = data.get("file")
        if file_field:
            abs_path = (asset_dir / file_field).resolve()
        else:
            # Try to match by stem
            abs_path = None
            for ext in _SUPPORTED_RASTER | {".svg", ".xml"}:
                candidate = asset_dir / f"{y.stem}{ext}"
                if candidate.exists():
                    abs_path = candidate.resolve()
                    break
        data["_abs_path"] = str(abs_path) if abs_path else None
        index[asset_id] = data
    return index


def _parse_placeholder_name(name: str) -> tuple[str, str] | None:
    if not name.startswith(PLACEHOLDER_PREFIX):
        return None
    rest = name[len(PLACEHOLDER_PREFIX):]
    parts = rest.split(":")
    asset_id = parts[0] if parts else None
    fit_mode = parts[1] if len(parts) > 1 else "fill"
    return asset_id, fit_mode


def _compute_crop_for_fill(slot_w: int, slot_h: int, img_w: int, img_h: int) -> dict:
    """Return python-pptx crop fractions (0..1) to center-crop image to slot."""
    if img_w == 0 or img_h == 0:
        return {"left": 0.0, "right": 0.0, "top": 0.0, "bottom": 0.0}
    slot_aspect = slot_w / slot_h
    img_aspect = img_w / img_h
    if img_aspect > slot_aspect:
        # Crop sides
        crop_total = 1 - slot_aspect / img_aspect
        side = crop_total / 2
        return {"left": side, "right": side, "top": 0.0, "bottom": 0.0}
    elif img_aspect < slot_aspect:
        # Crop top/bottom
        crop_total = 1 - img_aspect / slot_aspect
        side = crop_total / 2
        return {"left": 0.0, "right": 0.0, "top": side, "bottom": side}
    else:
        return {"left": 0.0, "right": 0.0, "top": 0.0, "bottom": 0.0}


def _splice_slide(slide, asset_index: dict, warnings: list) -> None:
    """Walk shapes in a slide and replace placeholders with pictures."""
    to_remove = []
    to_add = []

    for shape in slide.shapes:
        parsed = _parse_placeholder_name(shape.name or "")
        if not parsed:
            continue
        asset_id, fit_mode = parsed
        if asset_id == "none":
            continue  # legitimate placeholder, leave in place

        asset = asset_index.get(asset_id)
        if not asset or not asset.get("_abs_path"):
            warnings.append(
                f"missing asset id='{asset_id}' on slide {slide.slide_id}; "
                f"placeholder retained."
            )
            continue

        abs_path = Path(asset["_abs_path"])
        if abs_path.suffix.lower() not in _SUPPORTED_RASTER:
            warnings.append(
                f"unsupported asset format for '{asset_id}': {abs_path.suffix}; "
                f"placeholder retained."
            )
            continue

        left, top = shape.left, shape.top
        w, h = shape.width, shape.height
        to_remove.append(shape)
        to_add.append({
            "path": str(abs_path),
            "left": left, "top": top, "width": w, "height": h,
            "fit_mode": fit_mode,
        })

    # Remove placeholders
    for shape in to_remove:
        sp = shape._element
        sp.getparent().remove(sp)

    # Add pictures
    for spec in to_add:
        try:
            img = Image.open(spec["path"])
            img_w, img_h = img.size
        except Exception:
            img_w, img_h = (0, 0)

        if spec["fit_mode"] == "fill":
            pic = slide.shapes.add_picture(
                spec["path"], spec["left"], spec["top"],
                width=spec["width"], height=spec["height"],
            )
            crop = _compute_crop_for_fill(spec["width"], spec["height"], img_w, img_h)
            pic.crop_left = crop["left"]
            pic.crop_right = crop["right"]
            pic.crop_top = crop["top"]
            pic.crop_bottom = crop["bottom"]
        else:  # contain — letterbox
            if img_w and img_h:
                slot_aspect = spec["width"] / spec["height"]
                img_aspect = img_w / img_h
                if img_aspect > slot_aspect:
                    new_w = spec["width"]
                    new_h = int(new_w / img_aspect)
                else:
                    new_h = spec["height"]
                    new_w = int(new_h * img_aspect)
                offset_x = (spec["width"] - new_w) // 2
                offset_y = (spec["height"] - new_h) // 2
                slide.shapes.add_picture(
                    spec["path"],
                    spec["left"] + offset_x,
                    spec["top"] + offset_y,
                    width=new_w, height=new_h,
                )
            else:
                slide.shapes.add_picture(
                    spec["path"], spec["left"], spec["top"],
                    width=spec["width"], height=spec["height"],
                )


def splice(in_path: str, asset_dir: str, out_path: str) -> list[str]:
    prs = Presentation(in_path)
    asset_index = _read_asset_index(Path(asset_dir))
    warnings: list[str] = []
    for slide in prs.slides:
        _splice_slide(slide, asset_index, warnings)
    prs.save(out_path)

    if warnings:
        warn_path = out_path + ".warnings.txt"
        with open(warn_path, "w") as f:
            f.write("\n".join(warnings) + "\n")
    return warnings


def _cli():
    p = argparse.ArgumentParser(prog="splice_assets")
    p.add_argument("in_pptx")
    p.add_argument("--assets", required=True, help="asset directory with sidecar YAMLs")
    p.add_argument("-o", "--out", required=True, help="output .pptx path")
    args = p.parse_args()

    warnings = splice(args.in_pptx, args.assets, args.out)
    print(f"wrote {args.out}", file=sys.stderr)
    if warnings:
        print(f"{len(warnings)} warning(s) — see {args.out}.warnings.txt",
              file=sys.stderr)


if __name__ == "__main__":
    _cli()
