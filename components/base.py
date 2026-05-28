"""Component renderers.

A component is one piece of a slide (heading, text, image, metric, …) drawn at
a fractional rect. Each renderer is theme-aware via Context.

For image components, the renderer emits a PLACEHOLDER shape (dashed grey
rect labeled with the asset_id). splice_assets.py walks the final .pptx,
finds shapes whose name starts with `ASSET_PLACEHOLDER:`, and swaps in the
binary from the asset folder.
"""

from __future__ import annotations

from dataclasses import dataclass
from typing import Any

from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Inches, Pt


PLACEHOLDER_PREFIX = "ASSET_PLACEHOLDER"  # splice_assets.py looks for this


# ---------- Context + helpers ----------


@dataclass
class Context:
    """Render-time context passed to every component renderer."""

    theme: dict
    canvas_w_in: float
    canvas_h_in: float

    @classmethod
    def from_theme(cls, theme: dict) -> "Context":
        canvas = theme.get("canvas", {})
        return cls(
            theme=theme,
            canvas_w_in=canvas.get("width_in", 13.333),
            canvas_h_in=canvas.get("height_in", 7.5),
        )

    def to_emu(self, rect: dict) -> tuple[int, int, int, int]:
        """Fractional rect → (left, top, width, height) in EMU."""
        return (
            Inches(rect["x"] * self.canvas_w_in),
            Inches(rect["y"] * self.canvas_h_in),
            Inches(rect["w"] * self.canvas_w_in),
            Inches(rect["h"] * self.canvas_h_in),
        )

    def hex(self, color_key: str) -> str:
        """Resolve a color key to a hex string.

        Accepts:
          - direct hex: '#FFFFFF'
          - dotted: 'status.positive', 'tints.orange_30'
          - bare (searches palette → tints → status, first hit wins):
            'accent_primary', 'grey_90'
        """
        if color_key.startswith("#"):
            return color_key
        if "." in color_key:
            section, key = color_key.split(".", 1)
            return self.theme[section][key]
        for section in ("palette", "tints", "status"):
            if color_key in self.theme.get(section, {}):
                return self.theme[section][color_key]
        raise KeyError(f"color key not found in palette/tints/status: {color_key}")

    def rgb(self, color_key: str) -> RGBColor:
        return _hex_to_rgb(self.hex(color_key))

    def type_style(self, level: str) -> dict:
        return self.theme["type_scale"][level]

    def font_name(self, font_key: str) -> str:
        return self.theme["fonts"][font_key]


def _hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _apply_style_overrides(base: dict, overrides: dict | None) -> dict:
    if not overrides:
        return base
    out = dict(base)
    out.update(overrides)
    return out


def _set_run(run, *, text=None, font=None, size_pt=None, bold=None, italic=None,
             color_hex=None):
    if text is not None:
        run.text = text
    f = run.font
    if font is not None:
        f.name = font
    if size_pt is not None:
        f.size = Pt(size_pt)
    if bold is not None:
        f.bold = bold
    if italic is not None:
        f.italic = italic
    if color_hex is not None:
        f.color.rgb = _hex_to_rgb(color_hex)


def _add_textbox(slide, rect_emu, *, anchor=MSO_ANCHOR.TOP, word_wrap=True):
    left, top, w, h = rect_emu
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = word_wrap
    tf.vertical_anchor = anchor
    # Belt-and-braces: never auto-shrink or auto-grow. We size text explicitly
    # via measure_text and the type scale. PowerPoint sometimes enables
    # shrink-on-overflow by default for certain shape kinds.
    tf.auto_size = MSO_AUTO_SIZE.NONE
    # Clear default paragraph so we control insertion
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    return box, tf


_ALIGN = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}


# ---------- Component renderers ----------


_VERT_ANCHORS = {
    "top": MSO_ANCHOR.TOP,
    "middle": MSO_ANCHOR.MIDDLE,
    "bottom": MSO_ANCHOR.BOTTOM,
}


def render_heading(slide, ctx: Context, rect: dict, content: Any,
                   style_overrides: dict | None = None,
                   level: str = "h1", alignment: str = "left",
                   color_key: str = "text_primary",
                   vertical_anchor: str = None):
    """heading — content is a string. `level` keys into theme.type_scale.

    `vertical_anchor` overrides the per-level default (section_number → middle,
    others → top).
    """
    text = str(content) if content is not None else ""
    base = ctx.type_style(level)
    style = _apply_style_overrides(base, style_overrides)

    if vertical_anchor:
        anchor = _VERT_ANCHORS.get(vertical_anchor, MSO_ANCHOR.TOP)
    else:
        anchor = MSO_ANCHOR.MIDDLE if level == "section_number" else MSO_ANCHOR.TOP

    # Section numbers must never wrap — at 240pt even a 2-char string can
    # exceed any reasonable cell width.
    word_wrap = (level != "section_number")

    rect_emu = ctx.to_emu(rect)
    _, tf = _add_textbox(slide, rect_emu, anchor=anchor, word_wrap=word_wrap)

    p = tf.paragraphs[0]
    p.alignment = _ALIGN.get(alignment, PP_ALIGN.LEFT)
    run = p.add_run()
    _set_run(
        run,
        text=text,
        font=ctx.font_name(style["font"]),
        size_pt=style["size_pt"],
        bold=(style["weight"] >= 700),
        color_hex=style.get("color_hex") or ctx.hex(color_key),
    )


def render_text(slide, ctx: Context, rect: dict, content: Any,
                style_overrides: dict | None = None,
                level: str = "body", alignment: str = "left",
                color_key: str = "text_primary",
                vertical_anchor: str = "top"):
    """text — content is a string (paragraph) or list of strings (bullets).

    `vertical_anchor`: top (default) | middle | bottom — anchors text within
    its cell. Useful for centering a text block against an image of equal
    height.
    """
    base = ctx.type_style(level)
    style = _apply_style_overrides(base, style_overrides)
    rect_emu = ctx.to_emu(rect)
    anchor = _VERT_ANCHORS.get(vertical_anchor, MSO_ANCHOR.TOP)
    _, tf = _add_textbox(slide, rect_emu, anchor=anchor)

    font = ctx.font_name(style["font"])
    size_pt = style["size_pt"]
    bold = style["weight"] >= 700
    color = style.get("color_hex") or ctx.hex(color_key)

    if isinstance(content, list):
        # Bullets — each item is a paragraph
        for i, item in enumerate(content):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = _ALIGN.get(alignment, PP_ALIGN.LEFT)
            run = p.add_run()
            _set_run(
                run,
                text=f"•  {item}",
                font=font, size_pt=size_pt, bold=bold, color_hex=color,
            )
            p.space_after = Pt(4)
    else:
        text = str(content) if content is not None else ""
        p = tf.paragraphs[0]
        p.alignment = _ALIGN.get(alignment, PP_ALIGN.LEFT)
        run = p.add_run()
        _set_run(run, text=text, font=font, size_pt=size_pt, bold=bold,
                 color_hex=color)


def render_image(slide, ctx: Context, rect: dict, content: Any,
                 style_overrides: dict | None = None):
    """image — content is {"asset_id": "...", "fit": "fill"|"contain"}
    or {"placeholder": true, "label": "..."}.

    Always emits a placeholder shape. splice_assets.py walks the deck after
    render and swaps it for the binary.
    """
    rect_emu = ctx.to_emu(rect)
    left, top, w, h = rect_emu

    if isinstance(content, dict):
        asset_id = content.get("asset_id")
        fit = content.get("fit", "fill")
        label = content.get("label") or (asset_id or "image needed")
    else:
        asset_id = None
        fit = "fill"
        label = str(content) if content else "image needed"

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    # Marker for splice
    shape.name = f"{PLACEHOLDER_PREFIX}:{asset_id or 'none'}:{fit}"

    # Visual: dashed grey rect with label
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = ctx.rgb("tints.grey_60") if "tints" in ctx.theme else _hex_to_rgb("#EEEFF1")

    line = shape.line
    line.color.rgb = ctx.rgb("neutral_medium")
    line.width = Emu(12700)  # ~1pt
    line.dash_style = 7  # MSO_LINE_DASH_STYLE.DASH

    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    _set_run(
        run,
        text=label,
        font=ctx.font_name("body"),
        size_pt=11,
        color_hex=ctx.hex("text_secondary"),
    )


def render_metric(slide, ctx: Context, rect: dict, content: Any,
                  style_overrides: dict | None = None):
    """metric — content is {"value": "$1.8M", "label": "Revenue",
    "delta": "+12%"?, "delta_status": "positive"|"negative"|"warning"?}.

    Layout inside the cell: value (60%) top, label (25%) below, delta (15%) bottom.
    """
    rect_emu = ctx.to_emu(rect)
    left, top, w, h = rect_emu

    if not isinstance(content, dict):
        content = {"value": str(content), "label": ""}

    value = str(content.get("value", ""))
    label = str(content.get("label", ""))
    delta = content.get("delta")
    delta_status = content.get("delta_status", "positive")

    # Value
    value_h = int(h * 0.55)
    box = slide.shapes.add_textbox(left, top, w, value_h)
    tf = box.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    style = ctx.type_style("metric_value")
    run = p.add_run()
    _set_run(
        run,
        text=value,
        font=ctx.font_name(style["font"]),
        size_pt=style["size_pt"],
        bold=True,
        color_hex=ctx.hex("accent_primary"),
    )

    # Label
    label_top = top + value_h
    label_h = int(h * 0.25)
    box = slide.shapes.add_textbox(left, label_top, w, label_h)
    tf = box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    label_style = ctx.type_style("metric_label")
    run = p.add_run()
    _set_run(
        run,
        text=label,
        font=ctx.font_name(label_style["font"]),
        size_pt=label_style["size_pt"],
        color_hex=ctx.hex("text_secondary"),
    )

    # Delta
    if delta:
        delta_top = label_top + label_h
        delta_h = h - value_h - label_h
        box = slide.shapes.add_textbox(left, delta_top, w, delta_h)
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        _set_run(
            run,
            text=str(delta),
            font=ctx.font_name("body"),
            size_pt=14,
            bold=True,
            color_hex=ctx.hex(f"status.{delta_status}"),
        )


_CHART_TYPE_MAP = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "column_stacked": XL_CHART_TYPE.COLUMN_STACKED,
    "bar_stacked": XL_CHART_TYPE.BAR_STACKED,
    "line": XL_CHART_TYPE.LINE,
    "line_markers": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
    "area": XL_CHART_TYPE.AREA,
    "area_stacked": XL_CHART_TYPE.AREA_STACKED,
}


def render_chart(slide, ctx: Context, rect: dict, content: Any,
                 style_overrides: dict | None = None):
    """chart — content is {"type": "...", "categories": [...],
    "series": [{"name": "...", "values": [...]}, ...]}."""
    if not isinstance(content, dict):
        raise ValueError("chart content must be a dict with type, categories, series")

    rect_emu = ctx.to_emu(rect)
    left, top, w, h = rect_emu

    chart_type = _CHART_TYPE_MAP.get(content.get("type", "column"),
                                     XL_CHART_TYPE.COLUMN_CLUSTERED)
    cats = content.get("categories", [])
    series = content.get("series", [])

    data = CategoryChartData()
    data.categories = cats
    for s in series:
        data.add_series(s.get("name", ""), s.get("values", []))

    chart_shape = slide.shapes.add_chart(chart_type, left, top, w, h, data)
    chart = chart_shape.chart

    chart.has_title = False
    if len(series) > 1:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    else:
        chart.has_legend = False

    # Recolor series with theme accents
    accent_keys = ["accent_primary", "accent_secondary", "accent_tertiary",
                   "neutral_strong", "neutral_medium", "neutral_soft"]
    for i, plot_series in enumerate(chart.plots[0].series):
        color = ctx.rgb(accent_keys[i % len(accent_keys)])
        try:
            fill = plot_series.format.fill
            fill.solid()
            fill.fore_color.rgb = color
        except Exception:
            pass


TABLE_STYLES = {
    # header_accent: orange header band with white text. Body cells plain.
    # Thin orange line below header and between body rows.
    "header_accent": {
        "header_bg":   "accent_primary",
        "header_fg":   "text_inverse",
        "body_bg":     None,
        "body_fg":     "text_primary",
        "first_col_bg": None,
        "first_col_fg": None,
        "header_size": 14,
        "body_size":   12,
        "row_divider": "accent_primary",   # line below header AND every body row
    },
    # zebra_neutral: alternating light-grey body rows. Plain header with
    # orange underline to mark it. No row dividers (zebra does the separation).
    "zebra_neutral": {
        "header_bg":   None,
        "header_fg":   "text_primary",
        "body_bg_a":   "tints.grey_30",     # zebra-specific override
        "body_bg_b":   None,
        "body_fg":     "text_primary",
        "first_col_bg": None,
        "first_col_fg": None,
        "header_size": 14,
        "body_size":   12,
        "header_underline": "accent_primary",
    },
    # filled_accent: first column orange + black text; other body cells white
    # + black text. Orange row dividers between rows.
    "filled_accent": {
        "header_bg":   None,
        "header_fg":   "text_primary",
        "body_bg":     None,
        "body_fg":     "text_primary",
        "first_col_bg": "accent_primary",
        "first_col_fg": "text_primary",
        "header_size": 14,
        "body_size":   12,
        "row_divider": "accent_primary",
    },
    # filled_neutral: first column grey + black text; other body cells white
    # + black text. Grey row dividers between rows.
    "filled_neutral": {
        "header_bg":   None,
        "header_fg":   "text_primary",
        "body_bg":     None,
        "body_fg":     "text_primary",
        "first_col_bg": "tints.grey_30",
        "first_col_fg": "text_primary",
        "header_size": 14,
        "body_size":   12,
        "row_divider": "tints.grey_30",
    },
    # minimal: no fills. Orange underline below header only; no body dividers.
    "minimal": {
        "header_bg":   None,
        "header_fg":   "text_primary",
        "body_bg":     None,
        "body_fg":     "text_primary",
        "first_col_bg": None,
        "first_col_fg": None,
        "header_size": 14,
        "body_size":   12,
        "header_underline": "accent_primary",
    },
}


def _set_cell_bottom_border(cell, color_hex: str, width_pt: float = 1.0):
    """Add a bottom border to a cell via XML (python-pptx doesn't expose cell
    borders through high-level API)."""
    from lxml import etree
    from pptx.oxml.ns import qn

    tcPr = cell._tc.get_or_add_tcPr()
    # Remove any existing lnB
    for existing in tcPr.findall(qn("a:lnB")):
        tcPr.remove(existing)
    ln = etree.SubElement(tcPr, qn("a:lnB"))
    ln.set("w", str(int(width_pt * 12700)))
    ln.set("cap", "flat")
    ln.set("cmpd", "sng")
    ln.set("algn", "ctr")
    solid = etree.SubElement(ln, qn("a:solidFill"))
    color = etree.SubElement(solid, qn("a:srgbClr"))
    color.set("val", color_hex.lstrip("#").upper())
    prst = etree.SubElement(ln, qn("a:prstDash"))
    prst.set("val", "solid")


def render_table(slide, ctx: Context, rect: dict, content: Any,
                 style_overrides: dict | None = None):
    """table — content is {"rows": N, "cols": N, "has_header": bool,
    "data": [[...], ...], "style": "<style_name>"}.

    Available styles: header_accent (default), zebra_neutral, filled_accent,
    filled_neutral, minimal. See TABLE_STYLES for the parameters of each.

    The table shape is sized to fit the actual row content (compact rows),
    not stretched to fill the cell. Row heights:
      header: 0.45in   body: 0.40in
    Anything below the table within the cell area is left blank.
    """
    if not isinstance(content, dict):
        raise ValueError("table content must be a dict")

    rect_emu = ctx.to_emu(rect)
    left, top, slot_w, slot_h = rect_emu

    data = content.get("data", [])
    if not data:
        return
    rows = content.get("rows", len(data))
    cols = content.get("cols", max(len(r) for r in data) if data else 1)
    has_header = content.get("has_header", True)
    style_name = content.get("style", "header_accent")
    style = TABLE_STYLES.get(style_name, TABLE_STYLES["header_accent"])

    # Compact row heights — sized to fit text + small padding, not stretched
    # to fill the grid cell.
    header_h_in = 0.45
    body_h_in = 0.40
    if has_header:
        total_h_in = header_h_in + (rows - 1) * body_h_in
    else:
        total_h_in = rows * body_h_in
    table_h = Inches(total_h_in)

    table_shape = slide.shapes.add_table(rows, cols, left, top, slot_w, table_h)
    table = table_shape.table
    # Set explicit row heights (PowerPoint sometimes redistributes; this
    # ensures each row sticks to its target).
    for r in range(rows):
        is_header = has_header and r == 0
        table.rows[r].height = Inches(header_h_in if is_header else body_h_in)

    body_font = ctx.font_name("body")
    heading_font = ctx.font_name("heading")

    for r in range(rows):
        for c in range(cols):
            try:
                cell_text = str(data[r][c]) if r < len(data) and c < len(data[r]) else ""
            except (IndexError, TypeError):
                cell_text = ""
            cell = table.cell(r, c)
            cell.text = ""
            tf = cell.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            is_header = has_header and r == 0
            is_first_col = (c == 0) and not is_header

            # Choose fg / bg / size
            if is_header:
                fg_key = style["header_fg"]
                size_pt = style["header_size"]
                bg_key = style.get("header_bg")
                use_heading_font = True
            elif is_first_col and style.get("first_col_bg") is not None:
                fg_key = style.get("first_col_fg") or style["body_fg"]
                size_pt = style["body_size"]
                bg_key = style["first_col_bg"]
                use_heading_font = False
            else:
                fg_key = style["body_fg"]
                size_pt = style["body_size"]
                # zebra mode uses body_bg_a / body_bg_b alternating; other styles
                # use a single body_bg.
                if "body_bg_a" in style:
                    body_row_idx = r - (1 if has_header else 0)
                    bg_key = (style["body_bg_a"] if body_row_idx % 2 == 0
                              else style["body_bg_b"])
                else:
                    bg_key = style.get("body_bg")
                use_heading_font = False

            _set_run(
                run,
                text=cell_text,
                font=heading_font if use_heading_font else body_font,
                size_pt=size_pt,
                bold=is_header,
                color_hex=ctx.hex(fg_key),
            )

            fill = cell.fill
            if bg_key is None:
                fill.background()
            else:
                fill.solid()
                fill.fore_color.rgb = ctx.rgb(bg_key)

    # Borders: row dividers, header underline.
    divider_color = style.get("row_divider")
    header_underline = style.get("header_underline")
    for r in range(rows):
        is_header = has_header and r == 0
        is_last = (r == rows - 1)
        if is_header and header_underline:
            for c in range(cols):
                _set_cell_bottom_border(table.cell(r, c),
                                        ctx.hex(header_underline), width_pt=1.0)
        elif divider_color and not is_last:
            for c in range(cols):
                _set_cell_bottom_border(table.cell(r, c),
                                        ctx.hex(divider_color), width_pt=1.0)
        elif is_header and divider_color:
            for c in range(cols):
                _set_cell_bottom_border(table.cell(r, c),
                                        ctx.hex(divider_color), width_pt=1.0)


def render_quote(slide, ctx: Context, rect: dict, content: Any,
                 style_overrides: dict | None = None):
    """quote — content is {"text": "...", "attribution": "..."}."""
    if isinstance(content, str):
        content = {"text": content, "attribution": ""}
    text = content.get("text", "")
    attribution = content.get("attribution", "")

    rect_emu = ctx.to_emu(rect)
    left, top, w, h = rect_emu

    quote_h = int(h * 0.75)
    box = slide.shapes.add_textbox(left, top, w, quote_h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    style = ctx.type_style("quote")
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    _set_run(
        run,
        text=f"“{text}”",  # curly quotes
        font=ctx.font_name(style["font"]),
        size_pt=style["size_pt"],
        italic=True,
        color_hex=ctx.hex("text_primary"),
    )

    if attribution:
        attr_top = top + quote_h
        attr_h = h - quote_h
        box = slide.shapes.add_textbox(left, attr_top, w, attr_h)
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        _set_run(
            run,
            text=f"— {attribution}",
            font=ctx.font_name("body"),
            size_pt=14,
            color_hex=ctx.hex("text_secondary"),
        )


def render_cta(slide, ctx: Context, rect: dict, content: Any,
               style_overrides: dict | None = None):
    """cta — content is {"label": "...", "target": "...?"}.
    Renders as a filled accent rectangle with the label centered.
    """
    if isinstance(content, str):
        content = {"label": content}
    label = content.get("label", "")

    rect_emu = ctx.to_emu(rect)
    left, top, w, h = rect_emu

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ctx.rgb("accent_primary")
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    _set_run(
        run,
        text=label,
        font=ctx.font_name("heading"),
        size_pt=18,
        bold=True,
        color_hex=ctx.hex("text_inverse"),
    )


def render_divider(slide, ctx: Context, rect: dict, content: Any = None,
                   style_overrides: dict | None = None):
    """divider — thin horizontal rule. content is ignored."""
    rect_emu = ctx.to_emu(rect)
    left, top, w, h = rect_emu
    # Draw a thin line in the vertical middle of the rect
    line_h = Emu(12700)  # ~1pt
    mid = top + h // 2 - line_h // 2
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, mid, w, line_h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ctx.rgb("neutral_soft")
    shape.line.fill.background()


def render_spacer(slide, ctx: Context, rect: dict, content: Any = None,
                  style_overrides: dict | None = None):
    """spacer — explicit empty cell. Renders nothing."""
    return


def render_icon_label(slide, ctx: Context, rect: dict, content: Any,
                      style_overrides: dict | None = None):
    """icon_label — content is {"icon_asset_id": "...", "label": "..."}.

    Left ~25% is an icon placeholder; right ~75% is the label.
    """
    if isinstance(content, str):
        content = {"label": content}
    icon_id = content.get("icon_asset_id")
    label = content.get("label", "")

    rect_emu = ctx.to_emu(rect)
    left, top, w, h = rect_emu
    icon_w = int(w * 0.25)
    gap = int(w * 0.03)

    if icon_id or content.get("with_icon", True):
        # Icon placeholder
        icon_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, icon_w, h
        )
        icon_shape.name = f"{PLACEHOLDER_PREFIX}:{icon_id or 'none'}:contain"
        icon_shape.fill.solid()
        icon_shape.fill.fore_color.rgb = ctx.rgb("tints.orange_90") if "tints" in ctx.theme else _hex_to_rgb("#FFF5ED")
        icon_shape.line.color.rgb = ctx.rgb("accent_primary")
        icon_shape.line.width = Emu(12700)
        icon_shape.line.dash_style = 7

    label_left = left + icon_w + gap
    label_w = w - icon_w - gap
    box = slide.shapes.add_textbox(label_left, top, label_w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    _set_run(
        run,
        text=label,
        font=ctx.font_name("body"),
        size_pt=14,
        color_hex=ctx.hex("text_primary"),
    )


# ---------- Compound components ----------


def _normalize_image_content(image):
    """Accept asset_id string, dict with asset_id, or placeholder dict."""
    if image is None:
        return {"placeholder": True, "label": "image"}
    if isinstance(image, str):
        return {"asset_id": image, "fit": "fill"}
    if isinstance(image, dict):
        if "asset_id" in image or "placeholder" in image:
            return image
        return {"placeholder": True, "label": "image"}
    return {"placeholder": True, "label": "image"}


def render_card(slide, ctx: Context, rect: dict, content: Any,
                style_overrides: dict | None = None,
                variant: str = "image_left"):
    """Card — a compound component that lays out image + label + sublabel +
    body inside a single grid placement.

    The card handles its own internal proportions and gaps, so recipes that
    use cards (matrix_2x2, team_strip, team_grid_2x2, …) don't need to
    micro-place sub-components on the grid.

    content:
      image     ('asset_id' | {asset_id, fit} | {placeholder, label} | None)
      label     (str)         — main heading (h3)
      sublabel  (str)         — optional secondary line under label (caption,
                                secondary color). Use for role / subtitle.
      body      (str|list)    — paragraph or bullets

    variants:
      image_left  — image left ~30%, gap, text block right ~64%; text
                    vertically centered against the image
      image_top   — image top ~55%, gap, text block bottom ~40%
      text_only   — no image; label/sublabel top, body fills the rest

    Internal padding (gaps between sub-pieces) is built in. Recipes that
    place cards on the grid should add their own inter-card gaps via row /
    col gap cells.
    """
    if not isinstance(content, dict):
        return

    image = content.get("image")
    label = content.get("label", "")
    sublabel = content.get("sublabel", "")
    body = content.get("body", "")

    if variant == "image_top":
        _card_image_top(slide, ctx, rect, image, label, sublabel, body)
    elif variant == "text_only":
        _card_text_only(slide, ctx, rect, label, sublabel, body)
    else:
        _card_image_left(slide, ctx, rect, image, label, sublabel, body)


def _render_text_stack(slide, ctx, x, y, w, h, label, sublabel, body):
    """Render label + sublabel + body stacked vertically, top-anchored.

    Uses fixed proportions: label 25% of h, sublabel 18%, body 57%. If
    sublabel is empty, allocates its share to body.
    """
    has_sub = bool(sublabel)
    label_frac = 0.25
    sub_frac = 0.18 if has_sub else 0.0
    body_frac = 1.0 - label_frac - sub_frac

    label_h = h * label_frac
    sub_h = h * sub_frac
    body_h = h * body_frac

    cur_y = y
    if label:
        render_heading(slide, ctx,
                       {"x": x, "y": cur_y, "w": w, "h": label_h},
                       label, level="h3", vertical_anchor="top")
    cur_y += label_h
    if has_sub:
        render_text(slide, ctx,
                    {"x": x, "y": cur_y, "w": w, "h": sub_h},
                    sublabel, level="caption", color_key="text_secondary",
                    vertical_anchor="top")
        cur_y += sub_h
    if body:
        render_text(slide, ctx,
                    {"x": x, "y": cur_y, "w": w, "h": body_h},
                    body, level="body", vertical_anchor="top")


def _card_image_left(slide, ctx, rect, image, label, sublabel, body):
    # Fractions of card width: image 30%, gap 6%, text 64%.
    img_w = rect["w"] * 0.30
    gap_w = rect["w"] * 0.06
    txt_w = rect["w"] - img_w - gap_w
    txt_x = rect["x"] + img_w + gap_w

    render_image(
        slide, ctx,
        {"x": rect["x"], "y": rect["y"], "w": img_w, "h": rect["h"]},
        _normalize_image_content(image),
    )

    # Text stack centered vertically against the image.
    has_sub = bool(sublabel)
    block_frac = 0.85 if has_sub else 0.75
    block_h = rect["h"] * block_frac
    block_top = rect["y"] + (rect["h"] - block_h) / 2
    _render_text_stack(slide, ctx, txt_x, block_top, txt_w, block_h,
                       label, sublabel, body)


def _card_image_top(slide, ctx, rect, image, label, sublabel, body):
    img_h = rect["h"] * 0.55
    gap_h = rect["h"] * 0.05
    text_h = rect["h"] - img_h - gap_h

    render_image(
        slide, ctx,
        {"x": rect["x"], "y": rect["y"], "w": rect["w"], "h": img_h},
        _normalize_image_content(image),
    )
    _render_text_stack(
        slide, ctx, rect["x"], rect["y"] + img_h + gap_h,
        rect["w"], text_h, label, sublabel, body,
    )


def _card_text_only(slide, ctx, rect, label, sublabel, body):
    _render_text_stack(
        slide, ctx, rect["x"], rect["y"], rect["w"], rect["h"],
        label, sublabel, body,
    )


# ---------- Dispatch ----------


COMPONENT_RENDERERS = {
    "heading": render_heading,
    "text": render_text,
    "image": render_image,
    "metric": render_metric,
    "chart": render_chart,
    "table": render_table,
    "quote": render_quote,
    "cta": render_cta,
    "divider": render_divider,
    "spacer": render_spacer,
    "icon_label": render_icon_label,
    "card": render_card,
}


def render_component(slide, ctx: Context, placement: dict):
    """Top-level dispatch: take one component placement and render it.

    `placement` shape:
      {"type": "<component>", "grid": {...}, "content": ..., "style_overrides": {...}?, ...kwargs}

    Recipes emit lists of these; render.py loops over them per slide.
    """
    from grid import placement_to_rect

    comp_type = placement["type"]
    if comp_type not in COMPONENT_RENDERERS:
        raise ValueError(f"unknown component type: {comp_type}")

    rect = placement_to_rect(placement["grid"])
    extra = {k: v for k, v in placement.items()
             if k not in ("type", "grid", "content", "style_overrides")}
    COMPONENT_RENDERERS[comp_type](
        slide,
        ctx,
        rect,
        placement.get("content"),
        style_overrides=placement.get("style_overrides"),
        **extra,
    )
