"""render.py — turn a plan.json into a .pptx with image placeholders.

Usage:
  python render.py plan.json out.pptx [--theme theme.yaml]

After render, image components are placeholder shapes whose `name` attribute
starts with `ASSET_PLACEHOLDER:`. Run splice_assets.py to swap them for real
binaries.
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches

from components import Context, render_component


# Slide background enum → theme color key.
# white     → palette.background     (#FFFFFF)
# light_grey → palette.surface       (#EBEBEB)
# light_orange → tints.orange_60     (#FFE8D4)
BACKGROUND_COLOR_KEYS = {
    "white":        "background",
    "light_grey":   "surface",
    "light_orange": "tints.orange_60",
}


def _apply_background(slide, ctx: Context, bg_kind: str) -> None:
    """Paint a full-bleed rect as the slide background.

    python-pptx's slide.background only sets master-inherited fills, which
    don't always render in viewers; a full-canvas rect is more reliable.
    """
    if bg_kind in (None, "white"):
        return  # default white from the master
    color_key = BACKGROUND_COLOR_KEYS.get(bg_kind)
    if not color_key:
        raise ValueError(f"unknown background: {bg_kind}")
    canvas_w = Inches(ctx.canvas_w_in)
    canvas_h = Inches(ctx.canvas_h_in)
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, canvas_w, canvas_h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = ctx.rgb(color_key)
    rect.line.fill.background()
    rect.name = "SLIDE_BACKGROUND"
    # Send to back: move the bg rect XML to the start of spTree's shape list
    sp_tree = rect._element.getparent()
    sp_tree.remove(rect._element)
    # Find the first non-pictogram shape — insert after nvGrpSpPr/grpSpPr but
    # before other shapes. The simplest is to insert at index 2 (after the
    # two required header children).
    sp_tree.insert(2, rect._element)


from recipes import RECIPES


def render(plan: dict, theme: dict, out_path: str) -> None:
    prs = Presentation()
    canvas = theme.get("canvas", {})
    prs.slide_width = Inches(canvas.get("width_in", 13.333))
    prs.slide_height = Inches(canvas.get("height_in", 7.5))

    ctx = Context.from_theme(theme)

    # Use the blank layout — index 6 is typical for "Blank" in default master.
    blank_layout = prs.slide_layouts[6]

    for slide_spec in plan.get("slides", []):
        slide = prs.slides.add_slide(blank_layout)

        # Background first so it sits underneath everything.
        bg = slide_spec.get("background", "white")
        _apply_background(slide, ctx, bg)

        recipe = slide_spec.get("recipe")
        if recipe == "free":
            placements = slide_spec.get("components") or []
        else:
            if recipe not in RECIPES:
                raise ValueError(f"unknown recipe: {recipe}")
            placements = RECIPES[recipe](
                slide_spec.get("content") or {},
                **(slide_spec.get("params") or {}),
            )

        for placement in placements:
            render_component(slide, ctx, placement)

    prs.save(out_path)


def _load_yaml(path: str) -> dict:
    with open(path) as f:
        return yaml.safe_load(f)


def _load_json(path: str) -> dict:
    with open(path) as f:
        return json.load(f)


def _cli():
    p = argparse.ArgumentParser(prog="render")
    p.add_argument("plan", help="plan.json")
    p.add_argument("out", help="output .pptx path")
    p.add_argument("--theme", default=None,
                   help="theme.yaml path (defaults to ./theme.yaml)")
    args = p.parse_args()

    plan = _load_json(args.plan)
    theme_path = args.theme or str(Path(__file__).parent / "theme.yaml")
    theme = _load_yaml(theme_path)

    render(plan, theme, args.out)
    print(f"wrote {args.out}", file=sys.stderr)


if __name__ == "__main__":
    _cli()
