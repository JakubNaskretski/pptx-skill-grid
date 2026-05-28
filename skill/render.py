"""render.py — turn a plan.json into a .pptx.

Usage:
  python render.py plan.json out.pptx [--theme theme.yaml] [--assets DIR]
                                       [--no-splice]

Render produces a .pptx with image placeholders. If an `assets/` directory
exists (either passed via --assets or the bundled skill assets/ folder),
this script auto-splices the binaries in so the output is a ready-to-ship
deck. Pass --no-splice to keep the placeholder version (e.g. for re-splicing
later against a different asset folder).
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches

from components import Context, render_component


# --- ORG SETTINGS — replace with your values. Don't commit changes. ----------
# After editing, prevent accidental commits with:
#   git update-index --skip-worktree skill/render.py
# Reverse with: git update-index --no-skip-worktree skill/render.py
COMPANY_NAME = "company_name"           # ← your org name; appears in non-cover slide header
LOGO_ASSET_ID = "org_logo"              # ← skill/assets/<this>.{png,jpg,svg,webp}
PAGE_NUMBERS = True
SHOW_PRESENTATION_TITLE = True

# Where raster photo binaries live (outside the skill, too heavy to ship).
# None → splice uses its default ../assets-external/ (sibling of the repo).
# Set an absolute path if your external folder lives elsewhere.
EXTERNAL_ASSETS_DIR = None

# Pre-rendered opener slide. Drop your branded opening slide as a one-slide
# .pptx at skill/templates/opening-slide.pptx — it gets prepended to every
# deck (unless the agent's plan sets use_template_opener: false). See
# skill/templates/README.md for details.
USE_TEMPLATE_OPENER = True
TEMPLATE_OPENER_PATH = "templates/opening-slide.pptx"   # relative to skill/
# -----------------------------------------------------------------------------


# Slide background enum → theme color key.
# white     → palette.background     (#FFFFFF)
# light_grey → palette.surface       (#EBEBEB)
# light_orange → tints.orange_60     (#FFE8D4)
BACKGROUND_COLOR_KEYS = {
    "white":        "background",
    "light_grey":   "surface",
    "light_orange": "tints.orange_60",
}


def _apply_background(slide, ctx: Context, bg_kind: str) -> None:
    """Paint a full-bleed rect as the slide background.

    python-pptx's slide.background only sets master-inherited fills, which
    don't always render in viewers; a full-canvas rect is more reliable.
    """
    if bg_kind in (None, "white"):
        return  # default white from the master
    color_key = BACKGROUND_COLOR_KEYS.get(bg_kind)
    if not color_key:
        raise ValueError(f"unknown background: {bg_kind}")
    canvas_w = Inches(ctx.canvas_w_in)
    canvas_h = Inches(ctx.canvas_h_in)
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, canvas_w, canvas_h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = ctx.rgb(color_key)
    rect.line.fill.background()
    rect.name = "SLIDE_BACKGROUND"
    # Send to back: move the bg rect XML to the start of spTree's shape list
    sp_tree = rect._element.getparent()
    sp_tree.remove(rect._element)
    # Find the first non-pictogram shape — insert after nvGrpSpPr/grpSpPr but
    # before other shapes. The simplest is to insert at index 2 (after the
    # two required header children).
    sp_tree.insert(2, rect._element)


from recipes import RECIPES


def _apply_decorations(slide, ctx: Context, slide_spec: dict, plan: dict,
                       total_slides: int, is_cover: bool = False,
                       page_offset: int = 0) -> None:
    """Stamp logo, presentation title, and page number on every non-cover
    slide. Driven by the ORG SETTINGS constants at the top of this file.

    is_cover: skip decorations entirely (cover slide owns its own design).
    page_offset: added to slide_id when computing displayed page number
        (used to account for prepended template opener slides).
    """
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    if is_cover:
        return
    slide_id = slide_spec.get("id", 0)

    canvas_w_in = ctx.canvas_w_in
    canvas_h_in = ctx.canvas_h_in
    footer_y_in = canvas_h_in * 13 / 14 + 0.06  # just below content area

    # --- 1. Presentation title — footer strip, after the logo ---
    # Logo sits at x = canvas_w * 0.04 with width ≈ 0.40". This box starts
    # at canvas_w * 0.12 and runs up to ~canvas_w * 0.84 (page number sits
    # at 0.85+). Same y as the logo so they share one footer line.
    if SHOW_PRESENTATION_TITLE:
        # The literal placeholder "company_name" is intentional — skip the
        # text entirely until the user replaces it with their real org name.
        company = (COMPANY_NAME or "").strip()
        if company == "company_name":
            company = ""
        deck_title = (plan.get("deck_title") or "").strip()
        parts = [p for p in (company, deck_title) if p]
        if parts:
            text = " · ".join(parts)
            box = slide.shapes.add_textbox(
                Inches(canvas_w_in * 0.12),
                Inches(footer_y_in),
                Inches(canvas_w_in * 0.72),
                Inches(0.28),
            )
            tf = box.text_frame
            tf.word_wrap = False
            tf.margin_left = Emu(0)
            tf.margin_right = Emu(0)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = text
            run.font.name = ctx.font_name("body")
            run.font.size = Pt(10)
            run.font.color.rgb = ctx.rgb("text_secondary")

    # --- 2. Page number at footer-right ---
    if PAGE_NUMBERS:
        text = f"{slide_id + page_offset} / {total_slides}"
        box = slide.shapes.add_textbox(
            Inches(canvas_w_in * 0.85),
            Inches(footer_y_in),
            Inches(canvas_w_in * 0.11),
            Inches(0.28),
        )
        tf = box.text_frame
        tf.word_wrap = False
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = text
        run.font.name = ctx.font_name("body")
        run.font.size = Pt(10)
        run.font.color.rgb = ctx.rgb("text_secondary")

    # --- 3. Logo at footer-left ---
    if LOGO_ASSET_ID:
        asset_dir = Path(__file__).parent / "assets"
        logo_path = None
        for ext in (".png", ".jpg", ".jpeg", ".webp", ".svg"):
            cand = asset_dir / f"{LOGO_ASSET_ID}{ext}"
            if cand.exists():
                logo_path = cand
                break
        x = Inches(canvas_w_in * 0.04)
        y = Inches(footer_y_in - 0.05)
        h = Inches(0.40)
        if logo_path and logo_path.suffix.lower() == ".svg":
            # SVG logo: rasterize to a small high-res PNG and embed. We don't
            # use the native asvg:svgBlip path here because the logo is tiny
            # (~40px tall) and the raster is visually indistinguishable.
            try:
                import io
                import cairosvg  # type: ignore
                png_buf = io.BytesIO()
                cairosvg.svg2png(url=str(logo_path), output_width=1024,
                                 write_to=png_buf)
                png_buf.seek(0)
                slide.shapes.add_picture(png_buf, x, y, height=h)
            except ImportError:
                # cairosvg missing — show a hint-shaped placeholder.
                shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y,
                                               Inches(0.40), h)
                shape.name = f"LOGO_SVG_NEEDS_CAIROSVG:{LOGO_ASSET_ID}"
                shape.fill.solid()
                shape.fill.fore_color.rgb = ctx.rgb("tints.grey_60")
                shape.line.color.rgb = ctx.rgb("neutral_medium")
                shape.line.width = Emu(12700)
                shape.line.dash_style = 7
                tf = shape.text_frame
                p = tf.paragraphs[0]
                from pptx.enum.text import PP_ALIGN as _A
                p.alignment = _A.CENTER
                run = p.add_run()
                run.text = "svg"
                run.font.name = ctx.font_name("body")
                run.font.size = Pt(8)
                run.font.color.rgb = ctx.rgb("text_secondary")
        elif logo_path:
            slide.shapes.add_picture(str(logo_path), x, y, height=h)
        else:
            # Placeholder so the user sees "logo missing" until they add it
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y,
                                           Inches(0.40), h)
            shape.name = f"ASSET_PLACEHOLDER:{LOGO_ASSET_ID}:contain"
            shape.fill.solid()
            shape.fill.fore_color.rgb = ctx.rgb("tints.grey_60")
            shape.line.color.rgb = ctx.rgb("neutral_medium")
            shape.line.width = Emu(12700)
            shape.line.dash_style = 7
            tf = shape.text_frame
            tf.margin_left = Emu(0)
            tf.margin_right = Emu(0)
            p = tf.paragraphs[0]
            from pptx.enum.text import PP_ALIGN as _A
            p.alignment = _A.CENTER
            run = p.add_run()
            run.text = "logo"
            run.font.name = ctx.font_name("body")
            run.font.size = Pt(8)
            run.font.color.rgb = ctx.rgb("text_secondary")


def _resolve_template_opener(plan: dict, override_disabled: bool) -> Path | None:
    """Return the path to the template opener if it should be used, else None.

    Decision rules (in order):
      1. --no-template-opener CLI flag (override_disabled=True) → None
      2. USE_TEMPLATE_OPENER constant False → None
      3. plan.use_template_opener == False (explicit opt-out) → None
      4. Template file missing on disk → None (with a stderr warning)
      5. Otherwise → the resolved Path
    """
    if override_disabled or not USE_TEMPLATE_OPENER:
        return None
    if plan.get("use_template_opener") is False:
        return None
    template_path = Path(__file__).parent / TEMPLATE_OPENER_PATH
    if not template_path.exists():
        # Configured but missing — warn and proceed without it.
        print(f"(template opener configured but missing: {template_path})",
              file=sys.stderr)
        return None
    return template_path


def render(plan: dict, theme: dict, out_path: str,
           no_template_opener: bool = False) -> None:
    template_path = _resolve_template_opener(plan, no_template_opener)
    has_template = template_path is not None

    if has_template:
        # Load from template — its slide becomes slide 1, masters/layouts
        # come along. Agent's slides will append after.
        prs = Presentation(str(template_path))
    else:
        prs = Presentation()
        canvas = theme.get("canvas", {})
        prs.slide_width = Inches(canvas.get("width_in", 13.333))
        prs.slide_height = Inches(canvas.get("height_in", 7.5))

    ctx = Context.from_theme(theme)

    # Pick the blank layout. Standard PowerPoint masters expose 7 layouts
    # with Blank at index 6; if the template was saved from a non-standard
    # base it may have fewer — fall back to the first layout in that case.
    try:
        blank_layout = prs.slide_layouts[6]
    except IndexError:
        blank_layout = prs.slide_layouts[0]

    template_offset = 1 if has_template else 0
    total_slides = template_offset + len(plan.get("slides", []))

    for slide_spec in plan.get("slides", []):
        slide = prs.slides.add_slide(blank_layout)

        # Strip any placeholders the layout brought in. Even PowerPoint's
        # default "Blank" master/layout pair carries title + slide-number +
        # footer + section-index placeholders that bleed through as "Add
        # title" / "0" / etc. underneath our content. Killing them here
        # gives us a truly clean canvas regardless of which template (or
        # python-pptx default) the .pptx originated from.
        for ph in list(slide.placeholders):
            sp = ph._element
            sp.getparent().remove(sp)

        # Background first so it sits underneath everything.
        bg = slide_spec.get("background", "white")
        _apply_background(slide, ctx, bg)

        recipe = slide_spec.get("recipe")
        if recipe == "free":
            placements = slide_spec.get("components") or []
        else:
            if recipe not in RECIPES:
                raise ValueError(f"unknown recipe: {recipe}")
            placements = RECIPES[recipe](
                slide_spec.get("content") or {},
                **(slide_spec.get("params") or {}),
            )

        for placement in placements:
            render_component(slide, ctx, placement)

        # Decoration / cover rules:
        #   - With template: template owns the cover, so no agent slide is
        #     "the cover" — decorate every slide. Page numbers offset by 1.
        #   - Without template: agent's id=1 is the cover, skip decorations.
        is_cover = (not has_template) and slide_spec.get("id") == 1
        _apply_decorations(
            slide, ctx, slide_spec, plan, total_slides,
            is_cover=is_cover, page_offset=template_offset,
        )

    prs.save(out_path)


def _load_yaml(path: str) -> dict:
    with open(path) as f:
        return yaml.safe_load(f)


def _load_json(path: str) -> dict:
    with open(path) as f:
        return json.load(f)


DEFAULT_ASSETS_DIR = Path(__file__).parent / "assets"


def _should_splice(assets_dir: Path) -> bool:
    """True iff the directory has at least one sidecar .yaml that isn't the
    README/vocab/theme."""
    if not assets_dir.exists() or not assets_dir.is_dir():
        return False
    skip = {"theme.yaml", "asset_tag_vocab.yaml"}
    for p in assets_dir.glob("*.yaml"):
        if p.name not in skip:
            return True
    return False


def _cli():
    p = argparse.ArgumentParser(prog="render")
    p.add_argument("plan", help="plan.json")
    p.add_argument("out", help="output .pptx path")
    p.add_argument("--theme", default=None,
                   help="theme.yaml path (defaults to ./theme.yaml)")
    p.add_argument("--assets", default=None,
                   help="skill assets directory (defaults to bundled assets/). "
                        "Holds sidecar yamls + SVG binaries.")
    p.add_argument("--external-assets", default=None,
                   help="external raster folder (defaults to "
                        "EXTERNAL_ASSETS_DIR constant or ../assets-external/). "
                        "Holds photo binaries that live outside the skill.")
    p.add_argument("--no-splice", action="store_true",
                   help="Skip the splice step and leave image placeholders.")
    p.add_argument("--no-template-opener", action="store_true",
                   help="Skip the pre-rendered opening slide template for "
                        "this render, even if one is configured + on disk.")
    args = p.parse_args()

    plan = _load_json(args.plan)
    theme_path = args.theme or str(Path(__file__).parent / "theme.yaml")
    theme = _load_yaml(theme_path)

    render(plan, theme, args.out, no_template_opener=args.no_template_opener)
    print(f"wrote {args.out}", file=sys.stderr)

    # Auto-splice unless suppressed.
    if args.no_splice:
        return
    assets_dir = Path(args.assets) if args.assets else DEFAULT_ASSETS_DIR
    external_dir = args.external_assets or EXTERNAL_ASSETS_DIR
    if _should_splice(assets_dir):
        try:
            from splice_assets import splice
        except ImportError:
            print("(splice_assets.py not importable — skipping splice)", file=sys.stderr)
            return
        warnings = splice(
            args.out, args.out,
            skill_assets_dir=str(assets_dir),
            external_assets_dir=external_dir,
        )
        print(f"spliced from {assets_dir}"
              + (f" + external={external_dir}" if external_dir else "")
              + (f" ({len(warnings)} warnings)" if warnings else ""),
              file=sys.stderr)
    else:
        print(f"(no sidecars in {assets_dir} — leaving placeholders)",
              file=sys.stderr)


if __name__ == "__main__":
    _cli()
