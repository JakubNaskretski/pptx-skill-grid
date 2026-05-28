"""splice_assets.py — replace ASSET_PLACEHOLDER shapes with real images.

Walks a .pptx, finds each shape whose name starts with `ASSET_PLACEHOLDER:`,
parses the asset_id + fit_mode, looks up the binary, and inserts it at the
placeholder's geometry.

Asset binary resolution follows the two-folder layout:

  - .svg   → expected in skill/assets/   (lives with the yaml)
  - raster → expected in EXTERNAL_ASSETS_DIR (outside the skill; ../assets-external
             by default). Falls back to skill/assets/ for the logo case.

SVG handling: cairosvg rasterizes the SVG to a high-res PNG (2048px wide)
which is embedded via the standard picture path. Reliable in every viewer;
visually identical to vector at typical slide sizes.

Usage:
  python splice_assets.py in.pptx -o out.pptx
                          [--skill-assets DIR] [--external-assets DIR]
"""

from __future__ import annotations

import argparse
import io
import sys
from pathlib import Path

import yaml
from PIL import Image
from pptx import Presentation
from pptx.util import Emu

# cairosvg is only needed when there are SVG assets. Import is lazy so the
# splice path works for raster-only decks even when cairosvg isn't installed.
try:
    import cairosvg  # type: ignore
    _HAS_CAIROSVG = True
except ImportError:
    _HAS_CAIROSVG = False


PLACEHOLDER_PREFIX = "ASSET_PLACEHOLDER:"

_SUPPORTED_RASTER = {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff", ".webp"}
_SUPPORTED_VECTOR = {".svg"}

# Path resolution defaults
DEFAULT_SKILL_ASSETS = Path(__file__).parent / "assets"
DEFAULT_EXTERNAL_ASSETS = Path(__file__).parent.parent.parent / "assets-external"

# OOXML namespace (used only if we ever revisit native vector embed).
_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


# ---------- asset index ----------


def _read_asset_index(skill_assets_dir: Path) -> dict[str, dict]:
    """Build {asset_id → sidecar dict} from *.yaml files in skill_assets_dir.

    Does NOT pre-resolve binary paths — that's done at splice time per asset,
    based on file extension (SVG → skill, raster → external).
    """
    index: dict[str, dict] = {}
    for y in sorted(skill_assets_dir.glob("*.yaml")):
        if y.name in {"theme.yaml", "asset_tag_vocab.yaml"}:
            continue
        try:
            with open(y) as f:
                data = yaml.safe_load(f) or {}
        except yaml.YAMLError:
            continue
        asset_id = data.get("id") or y.stem
        index[asset_id] = data
    return index


def _resolve_binary_path(
    asset: dict,
    skill_assets_dir: Path,
    external_dir: Path | None,
) -> Path | None:
    """Locate the binary for an asset entry.

    Lookup order:
      - .svg  → skill_assets_dir/<file>
      - raster → external_dir/<file>, falling back to skill_assets_dir/<file>
                 (the fallback covers the logo case, where a raster lives
                  inside the skill rather than the external folder).
    """
    file = asset.get("file")
    if not file:
        return None
    ext = Path(file).suffix.lower()
    skill_path = skill_assets_dir / file
    if ext in _SUPPORTED_VECTOR:
        return skill_path if skill_path.exists() else None
    if external_dir:
        ext_path = external_dir / file
        if ext_path.exists():
            return ext_path
    return skill_path if skill_path.exists() else None


# ---------- placeholder parsing ----------


def _parse_placeholder_name(name: str) -> tuple[str, str] | None:
    if not name.startswith(PLACEHOLDER_PREFIX):
        return None
    rest = name[len(PLACEHOLDER_PREFIX):]
    parts = rest.split(":")
    asset_id = parts[0] if parts else None
    fit_mode = parts[1] if len(parts) > 1 else "fill"
    return asset_id, fit_mode


# ---------- geometry helpers ----------


def _compute_crop_for_fill(slot_w: int, slot_h: int, img_w: int, img_h: int) -> dict:
    """Crop fractions (0..1) that center-crop an image to fill a slot."""
    if img_w == 0 or img_h == 0:
        return {"left": 0.0, "right": 0.0, "top": 0.0, "bottom": 0.0}
    slot_aspect = slot_w / slot_h
    img_aspect = img_w / img_h
    if img_aspect > slot_aspect:
        crop_total = 1 - slot_aspect / img_aspect
        side = crop_total / 2
        return {"left": side, "right": side, "top": 0.0, "bottom": 0.0}
    if img_aspect < slot_aspect:
        crop_total = 1 - img_aspect / slot_aspect
        side = crop_total / 2
        return {"left": 0.0, "right": 0.0, "top": side, "bottom": side}
    return {"left": 0.0, "right": 0.0, "top": 0.0, "bottom": 0.0}


def _letterbox_rect(slot_w: int, slot_h: int, img_w: int, img_h: int) -> tuple[int, int, int, int]:
    """Center-letterbox: return (offset_x, offset_y, new_w, new_h) so an
    image of (img_w, img_h) fits inside (slot_w, slot_h) at preserved aspect.
    """
    if not (img_w and img_h):
        return 0, 0, slot_w, slot_h
    slot_aspect = slot_w / slot_h
    img_aspect = img_w / img_h
    if img_aspect > slot_aspect:
        new_w = slot_w
        new_h = int(slot_w / img_aspect)
    else:
        new_h = slot_h
        new_w = int(slot_h * img_aspect)
    return (slot_w - new_w) // 2, (slot_h - new_h) // 2, new_w, new_h


# ---------- SVG: native vector embed ----------


def _rasterize_svg_to_png(svg_path: Path, target_w_px: int = 2048) -> bytes:
    """Render an SVG to PNG bytes for the embedded fallback.

    The fallback is what PowerPoint <2016 (and any viewer that doesn't
    recognize the asvg:svgBlip extension) will display. 2048px wide is
    plenty for any reasonable zoom.
    """
    if not _HAS_CAIROSVG:
        raise RuntimeError(
            "SVG embed requires cairosvg. Install with: pip install cairosvg"
        )
    buf = io.BytesIO()
    cairosvg.svg2png(url=str(svg_path), output_width=target_w_px, write_to=buf)
    return buf.getvalue()


def _embed_svg(slide, svg_path: Path, left, top, width, height):
    """Rasterize an SVG to high-res PNG and embed it via add_picture.

    Returns the python-pptx Picture shape.

    Earlier versions of this function ALSO embedded the original SVG
    bytes as an asvg:svgBlip extension so PowerPoint 2016+ would render
    the picture as native vector. That path proved fragile across
    python-pptx versions — the ImagePart + relate_to dance for a
    non-raster content type, plus the namespace serialization for the
    extension XML, occasionally produced .pptx files PowerPoint refused
    to open. The rasterized PNG looks identical at typical slide sizes
    (we render at 2048px wide) and is reliable in every viewer.
    """
    png_bytes = _rasterize_svg_to_png(svg_path)
    return slide.shapes.add_picture(
        io.BytesIO(png_bytes), left, top, width=width, height=height,
    )


# ---------- per-slide splice ----------


def _splice_slide(
    slide,
    asset_index: dict,
    skill_assets_dir: Path,
    external_dir: Path | None,
    warnings: list,
) -> None:
    """Walk shapes in a slide and replace placeholders with pictures."""
    to_remove = []
    to_add = []

    for shape in slide.shapes:
        parsed = _parse_placeholder_name(shape.name or "")
        if not parsed:
            continue
        asset_id, fit_mode = parsed
        if asset_id == "none":
            continue  # legitimate "no asset" placeholder, leave in place

        asset = asset_index.get(asset_id)
        if not asset:
            warnings.append(
                f"missing asset id='{asset_id}' on slide {slide.slide_id}; "
                f"placeholder retained."
            )
            continue

        abs_path = _resolve_binary_path(asset, skill_assets_dir, external_dir)
        if not abs_path:
            warnings.append(
                f"binary not found for '{asset_id}' (file={asset.get('file')}, "
                f"checked skill_assets={skill_assets_dir} and "
                f"external={external_dir}); placeholder retained."
            )
            continue

        suffix = abs_path.suffix.lower()
        if suffix in _SUPPORTED_VECTOR and not _HAS_CAIROSVG:
            warnings.append(
                f"SVG asset '{asset_id}' requires cairosvg "
                f"(pip install cairosvg); placeholder retained."
            )
            continue
        if suffix not in _SUPPORTED_RASTER and suffix not in _SUPPORTED_VECTOR:
            warnings.append(
                f"unsupported format for '{asset_id}': {suffix}; "
                f"placeholder retained."
            )
            continue

        to_remove.append(shape)
        to_add.append({
            "asset": asset,
            "path": abs_path,
            "left": shape.left, "top": shape.top,
            "width": shape.width, "height": shape.height,
            "fit_mode": fit_mode,
        })

    for shape in to_remove:
        sp = shape._element
        sp.getparent().remove(sp)

    for spec in to_add:
        _insert_asset(slide, spec)


def _insert_asset(slide, spec: dict) -> None:
    """Place one asset at the spec's geometry, applying fit_mode."""
    abs_path = spec["path"]
    suffix = abs_path.suffix.lower()
    is_svg = suffix in _SUPPORTED_VECTOR
    left, top = spec["left"], spec["top"]
    w, h = spec["width"], spec["height"]
    fit_mode = spec["fit_mode"]

    # Intrinsic dims: from yaml for SVG (PIL can't read it), from PIL for raster.
    if is_svg:
        asset = spec["asset"]
        img_w = asset.get("width") or 0
        img_h = asset.get("height") or 0
    else:
        try:
            with Image.open(abs_path) as img:
                img_w, img_h = img.size
        except Exception:
            img_w, img_h = 0, 0

    if fit_mode == "fill":
        if is_svg:
            pic = _embed_svg(slide, abs_path, left, top, w, h)
        else:
            pic = slide.shapes.add_picture(
                str(abs_path), left, top, width=w, height=h,
            )
        crop = _compute_crop_for_fill(w, h, img_w, img_h)
        pic.crop_left = crop["left"]
        pic.crop_right = crop["right"]
        pic.crop_top = crop["top"]
        pic.crop_bottom = crop["bottom"]
        return

    # "contain" — letterbox preserving aspect, centered in slot.
    ox, oy, new_w, new_h = _letterbox_rect(w, h, img_w, img_h)
    if is_svg:
        _embed_svg(slide, abs_path, left + ox, top + oy, new_w, new_h)
    else:
        slide.shapes.add_picture(
            str(abs_path), left + ox, top + oy, width=new_w, height=new_h,
        )


# ---------- top-level splice ----------


def splice(
    in_path: str,
    out_path: str,
    skill_assets_dir: str | Path | None = None,
    external_assets_dir: str | Path | None = None,
) -> list[str]:
    """Replace every ASSET_PLACEHOLDER in in_path with real binaries.

    skill_assets_dir defaults to the bundled assets/ folder (where SVGs and
    sidecar yamls live). external_assets_dir defaults to ../assets-external
    (where raster photos live, outside the skill).
    """
    skill_assets_dir = Path(skill_assets_dir or DEFAULT_SKILL_ASSETS)
    external_dir = Path(external_assets_dir) if external_assets_dir else DEFAULT_EXTERNAL_ASSETS
    if not external_dir.exists():
        external_dir = None  # raster lookups will fall back to skill_assets_dir

    prs = Presentation(in_path)
    asset_index = _read_asset_index(skill_assets_dir)
    warnings: list[str] = []
    for slide in prs.slides:
        _splice_slide(slide, asset_index, skill_assets_dir, external_dir, warnings)
    prs.save(out_path)

    if warnings:
        warn_path = out_path + ".warnings.txt"
        with open(warn_path, "w") as f:
            f.write("\n".join(warnings) + "\n")
    return warnings


# ---------- CLI ----------


def _cli():
    p = argparse.ArgumentParser(prog="splice_assets")
    p.add_argument("in_pptx")
    p.add_argument("-o", "--out", required=True, help="output .pptx path")
    p.add_argument("--skill-assets", default=None,
                   help="dir containing sidecar yamls + SVG binaries "
                        "(default: bundled assets/)")
    p.add_argument("--external-assets", default=None,
                   help="dir containing raster photo binaries "
                        "(default: ../assets-external/)")
    args = p.parse_args()

    warnings = splice(
        args.in_pptx, args.out,
        skill_assets_dir=args.skill_assets,
        external_assets_dir=args.external_assets,
    )
    print(f"wrote {args.out}", file=sys.stderr)
    if warnings:
        print(f"{len(warnings)} warning(s) — see {args.out}.warnings.txt",
              file=sys.stderr)


if __name__ == "__main__":
    _cli()
